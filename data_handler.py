import pandas as pd
from pptx import Presentation
import os
from config import EXCEL_FILE_PATH, PPTX_FILE_PATH

def read_excel_data():
    """지정된 경로의 Excel 파일을 읽어 데이터와 경로를 반환합니다."""
    if not os.path.exists(EXCEL_FILE_PATH):
        return None, f"오류: Excel 파일 경로를 찾을 수 없습니다.\n경로: {EXCEL_FILE_PATH}"

    try:
        # pandas를 사용하여 Excel 파일을 읽습니다.
        df = pd.read_excel(EXCEL_FILE_PATH)
        output_content = df.to_string(index=False) 
        return output_content, EXCEL_FILE_PATH
    except Exception as e:
        return None, f"오류: Excel 파일 처리 중 오류 발생:\n{e}"

def read_pptx_data(menu_letter='K', item_label='ECID'):
    """지정된 경로의 PPTX 파일을 읽어 텍스트 데이터와 경로를 반환합니다."""
    if not os.path.exists(PPTX_FILE_PATH):
        return None, f"오류: PPTX 파일 경로를 찾을 수 없습니다.\n경로: {PPTX_FILE_PATH}"

    try:
        prs = Presentation(PPTX_FILE_PATH)
        extracted_text = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            if slide_text:
                extracted_text.append(f"--- 슬라이드 {slide_num + 1} ---\n" + "\n".join(slide_text))

        output_content = "\n\n".join(extracted_text)
        return output_content, PPTX_FILE_PATH
        
    except Exception as e:
        return None, f"오류: PPTX 파일 처리 중 오류 발생:\n{e}"
